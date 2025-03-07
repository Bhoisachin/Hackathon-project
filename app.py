import streamlit as st
from PyPDF2 import PdfReader
import os
import pytesseract
from PIL import Image
import cv2
import numpy as np
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

# Configure Google AI
import google.generativeai as genai
genai.configure(api_key=GOOGLE_API_KEY)

# Windows: Set Tesseract OCR Path (if not in PATH)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Function to extract text from uploaded files
def extract_text(files):
    text = ""
    
    for file in files:
        file_path = file.name

        # Determine file type and process accordingly
        if file_path.endswith(".pdf"):
            text += extract_text_from_pdf(file)
        elif file_path.endswith(".docx"):
            text += extract_text_from_docx(file)
        elif file_path.endswith(".pptx"):
            text += extract_text_from_pptx(file)
        elif file_path.endswith((".png", ".jpg", ".jpeg")):
            text += extract_text_from_image(file)

    return text

def extract_text_from_pdf(file):
    """ Extract text from PDFs using PyPDF2 and OCR for scanned PDFs. """
    text = ""
    pdf_reader = PdfReader(file)
    for page in pdf_reader.pages:
        text += page.extract_text() or ""

    if len(text.strip()) < 10:  # If no text, use OCR
        text += extract_text_from_images(file)

    return text

def extract_text_from_images(pdf_file):
    """ Convert scanned PDF pages to images and extract text using OCR. """
    images = convert_from_path(pdf_file)
    extracted_text = ""

    for img in images:
        img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2GRAY)
        _, img_cv = cv2.threshold(img_cv, 150, 255, cv2.THRESH_BINARY)
        img_pil = Image.fromarray(img_cv)
        extracted_text += pytesseract.image_to_string(img_pil) + "\n"

    return extracted_text

def extract_text_from_docx(file):
    """ Extract text from a DOCX (Word) file. """
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    """ Extract text from a PPTX (PowerPoint) file. """
    presentation = Presentation(file)
    text = ""
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def extract_text_from_image(file):
    """ Extract text from an image using OCR. """
    img = Image.open(file)
    return pytesseract.image_to_string(img)

# Function to process text and generate responses
def get_answer(question, text):
    """ Generate answers using Google Gemini AI """
    if not text.strip():
        return "No text found in uploaded files. Please upload a valid document."

    # Split text into chunks
    text_chunks = [text[i:i+10000] for i in range(0, len(text), 10000)]

    # Create FAISS vector store
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)

    # Retrieve relevant text from FAISS
    docs = vector_store.similarity_search(question)

    # Create conversational chain
    prompt_template = """
    Answer the question based on the provided context. 
    If the answer is not in the context, say "Answer is not available."
    
    Context:\n{context}\n
    Question:\n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    # Generate answer
    response = chain({"input_documents": docs, "question": question}, return_only_outputs=True)
    return response["output_text"]

# Streamlit app main function
def main():
    st.set_page_config(page_title="Chat with Documents", layout="wide")
    st.header("Chat with PDF, DOCX, PPTX & Images using Gemini 💁")

    # File uploader
    uploaded_files = st.file_uploader("Upload Files (PDF, DOCX, PPTX, Images)", 
                                      type=["pdf", "docx", "pptx", "png", "jpg", "jpeg"], 
                                      accept_multiple_files=True)

    if uploaded_files:
        with st.spinner("Extracting text..."):
            document_text = extract_text(uploaded_files)
            st.success("Text extracted successfully!")

        # User question input
        st.write("### Ask a Question:")
        user_question = st.text_input("Enter your question:")

        if st.button("Get Answer"):
            with st.spinner("Generating response..."):
                response = get_answer(user_question, document_text)
                st.write("### Reply:", response)

if __name__ == "__main__":
    main()
